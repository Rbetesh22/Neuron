from pathlib import Path
from .base import Document, _h


class FileIngester:
    def ingest(self, path: str) -> list[Document]:
        p = Path(path).expanduser().resolve()
        if not p.exists():
            raise FileNotFoundError(f"File not found: {path}")

        suffix = p.suffix.lower()
        if suffix == ".pdf":
            content = self._read_pdf(p)
        elif suffix in (".docx", ".doc"):
            content = self._read_docx(p)
        elif suffix == ".pptx":
            content = self._read_pptx(p)
        else:
            content = p.read_text(encoding="utf-8", errors="ignore")

        return [Document(
            id=f"file_{_h(str(p))}",
            content=content,
            source="file",
            title=p.name,
            metadata={"path": str(p), "type": suffix.lstrip(".")},
        )]

    def _read_pdf(self, path: Path) -> str:
        import logging
        import pypdf
        logging.getLogger("pypdf").setLevel(logging.ERROR)
        reader = pypdf.PdfReader(str(path), strict=False)
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    def _read_docx(self, path: Path) -> str:
        import docx
        doc = docx.Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    def _read_pptx(self, path: Path) -> str:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = para.text.strip()
                        if line:
                            texts.append(line)
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
